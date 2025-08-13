import os
import json
from PyPDF2 import PdfReader
from pptx import Presentation
from openai import OpenAI

# 初始化OpenAI客户端
client = OpenAI(api_key="您的API_KEY")

def extract_pdf_text(pdf_path):
    """
    提取PDF文本内容（带页码标记）
    
    参数:
        pdf_path: PDF文件路径
        
    返回:
        str: 拼接后的文本内容（前15000字符）
    """
    text = ""
    with open(pdf_path, 'rb') as file:
        reader = PdfReader(file)
        for i, page in enumerate(reader.pages):
            text += f"【Page {i+1}】{page.extract_text()}\n\n"
    return text[:15000]  # 控制处理长度

def identify_document_type(pdf_text):
    """
    自动识别PDF文档类型
    
    参数:
        pdf_text: PDF提取的文本
        
    返回:
        str: 文档类型（行业报告/企业财报/专业论文）
    """
    doc_types = ["行业报告", "企业财报", "专业论文"]
    
    prompt = f"""
    请判断以下经济学文档的类型（仅返回类型名称）：
    可选类型：{doc_types}
    
    文档开头内容：
    {pdf_text[:2000]}
    """
    
    response = client.chat.completions.create(
        model="gpt-4-turbo",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.1
    )
    
    doc_type = response.choices[0].message.content.strip()
    return doc_type if doc_type in doc_types else "行业报告"  # 默认值

def generate_type_specific_prompt(doc_type):
    """
    根据文档类型生成专属提示词模板
    
    参数:
        doc_type: 已识别的文档类型
        
    返回:
        str: 格式化后的多轮提示词
    """
    prompt_templates = {
        "行业报告": {
            "focus": "重点关注：市场增长率/竞争格局/政策影响",
            "charts": "推荐图表：趋势线图/市场份额饼图/区域热力图"
        },
        "企业财报": {
            "focus": "重点关注：营收构成/成本结构/现金流",
            "charts": "推荐图表：柱状对比图/财务指标雷达图/现金流桑基图"
        },
        "专业论文": {
            "focus": "重点关注：理论模型/实证结果/政策建议",
            "charts": "推荐图表：回归分析表/机制流程图/假设验证矩阵"
        }
    }
    
    return f"""
    作为经济学讲师，请为[{doc_type}]创建教学Keynote，要求：
    
    1. 内容结构：
       - 标题：突出核心结论（<12字）
       - 摘要页：3个数据支撑的发现
       - 分析页：按"背景-方法-结论"展开
       - 图表页：{prompt_templates[doc_type]['charts']}
    
    2. 专业要求：
       - 数据必须标注来源（如Table 3/P15）
       - 复杂概念添加【教学提示】注释
       - {prompt_templates[doc_type]['focus']}
    
    3. 输出格式：
    {{
      "title": "标题",
      "slides": [
        {{
          "title": "页面标题",
          "content": ["要点1", "要点2"],
          "chart_instruction": "图表描述"
        }},
        ... // 至少6页
      ]
    }}
    """

def analyze_pdf_content(pdf_text, doc_type):
    """
    分析PDF内容并生成Keynote结构
    
    参数:
        pdf_text: PDF文本内容
        doc_type: 文档类型
        
    返回:
        dict: Keynote结构数据
    """
    prompt = generate_type_specific_prompt(doc_type)
    
    full_prompt = f"""
    文档类型：{doc_type}
    文档内容：
    {pdf_text}
    
    生成要求：
    {prompt}
    """
    
    response = client.chat.completions.create(
        model="gpt-4-turbo",
        messages=[
            {"role": "system", "content": "您是为经济学教授生成Keynote的专家"},
            {"role": "user", "content": full_prompt}
        ],
        response_format={"type": "json_object"},
        temperature=0.3
    )
    
    return json.loads(response.choices[0].message.content)

def create_keynote(analysis_result, output_path="output.pptx"):
    """
    根据分析结果生成PPTX文件
    
    参数:
        analysis_result: 分析结果字典
        output_path: 输出文件路径
        
    返回:
        str: 生成的PPTX文件路径
    """
    prs = Presentation()
    
    # 封面页
    slide_layout = prs.slide_layouts[0]  # 标题版式
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = analysis_result["title"]
    
    # 内容页
    for slide_data in analysis_result["slides"]:
        slide_layout = prs.slide_layouts[1]  # 标题+内容版式
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title = slide.shapes.title
        title.text = slide_data["title"]
        
        # 添加内容
        content = slide.shapes.placeholders[1]
        tf = content.text_frame
        
        for point in slide_data["content"]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
        
        # 添加图表提示（备注）
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = f"图表建议：{slide_data.get('chart_instruction', '')}"
    
    prs.save(output_path)
    return os.path.abspath(output_path)

def process_pdf_to_keynote(pdf_path):
    """
    完整处理流程
    
    参数:
        pdf_path: 输入的PDF文件路径
        
    返回:
        tuple: (Keynote结构数据, 生成的PPTX路径)
    """
    print(f"开始处理: {os.path.basename(pdf_path)}")
    
    # 第一步：提取文本
    print("步骤1/4: 提取PDF文本...")
    pdf_text = extract_pdf_text(pdf_path)
    
    # 第二步：识别类型
    print("步骤2/4: 识别文档类型...")
    doc_type = identify_document_type(pdf_text)
    print(f"识别结果: {doc_type}")
    
    # 第三步：分析内容
    print("步骤3/4: 分析文档内容...")
    analysis = analyze_pdf_content(pdf_text, doc_type)
    
    # 第四步：生成Keynote
    print("步骤4/4: 生成PPTX文件...")
    output_path = create_keynote(analysis)
    
    return analysis, output_path

# 使用示例
if __name__ == "__main__":
    # 输入PDF路径
    input_pdf = "经济学论文.pdf"
    
    # 执行转换
    analysis_data, keynote_file = process_pdf_to_keynote(input_pdf)
    
    # 打印结果
    print("\n生成结果:")
    print(f"- Keynote结构: {json.dumps(analysis_data, indent=2, ensure_ascii=False)}")
    print(f"- 文件已保存到: {keynote_file}")